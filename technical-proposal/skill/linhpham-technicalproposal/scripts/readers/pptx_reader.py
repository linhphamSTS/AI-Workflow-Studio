#!/usr/bin/env python3
"""Read a .pptx deck and emit one Markdown section per slide so the
Phase-1 ingest agent can analyse architecture slides, kick-off decks, etc.

Captures: title, body text, speaker notes, image alt text where present.
Embedded images themselves are not extracted (use a separate image-OCR
pass if the deck is mostly screenshots)."""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def _ensure_pptx():
    try:
        import pptx  # noqa: F401
    except ImportError:
        import subprocess
        subprocess.run(
            [sys.executable, "-m", "pip", "install", "--quiet", "--user", "python-pptx"],
            check=True,
        )


def read(path: Path) -> str:
    _ensure_pptx()
    from pptx import Presentation
    prs = Presentation(str(path))
    out = [f"# Deck: {path.name}", ""]
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        out.append(f"- {text}")
            if shape.shape_type == 13 and getattr(shape, "image", None):
                # picture; capture alt text if any
                alt = getattr(shape, "alt_text", "") or ""
                if alt:
                    out.append(f"  ![image]({alt})")
        # speaker notes
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        if notes:
            out.append(f"\n**Speaker notes:** {notes}")
        out.append("")
    return "\n".join(out)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("path", type=Path)
    args = ap.parse_args()
    if not args.path.exists():
        print(f"! file not found: {args.path}", file=sys.stderr)
        return 1
    sys.stdout.write(read(args.path))
    return 0


if __name__ == "__main__":
    sys.exit(main())
